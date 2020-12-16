from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.oxml import parse_xml
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import xmltodict, json
from pptx.enum.dml import MSO_THEME_COLOR, MSO_COLOR_TYPE

def get_theme_info(slide_layout):
    # return parse_xml(slide_layout.slide_master.part.part_related_by(RT.THEME).blob)
    return xmltodict.parse(slide_layout.slide_master.part.part_related_by(RT.THEME).blob)

def get_info_from_shape(shape):
    # x = parse_xml(shape._element.xml)
    x = xmltodict.parse(shape._element.xml)
    try:
        defRPr = x['p:sp']['p:txBody']['a:lstStyle']['a:lvl1pPr']['a:defRPr']
        text = x['p:sp']['p:txBody']['a:p']['a:r']['a:t']
        font_bold = defRPr['@b'] if '@b' in defRPr else 0
        font_italic = defRPr['@i'] if '@i' in defRPr else 0
        font_size = int(defRPr['@sz'])/100 if '@sz' in defRPr else 0
        if 'a:latin' in defRPr:
            font_name = defRPr['a:latin']['@typeface'] if 'a:latin' in defRPr else ''
        else:
            font_name = theme_fonts['a:majorFont']['a:latin']['@typeface']

        if 'a:solidFill' in defRPr:
            if 'a:srgbClr' in defRPr['a:solidFill']:
                font_color = defRPr['a:solidFill']['a:srgbClr']['@val']
                font_type = MSO_COLOR_TYPE.RGB
            elif 'a:schemeClr' in defRPr['a:solidFill']:
                font_scheme_color = defRPr['a:solidFill']['a:schemeClr']['@val']
                if f"a:{font_scheme_color}" in theme_colors:
                    font_color = theme_colors[f"a:{font_scheme_color}"]['a:srgbClr']['@val']
                else:
                    font_color = default_clr
                font_type = MSO_COLOR_TYPE.SCHEME
            else:
                font_color = ''
                font_type = ''
        else:
            font_color = ''
            font_type = ''

        return text, font_bold, font_italic, font_name, font_size, font_type, font_color
    except:
        return [""]*7

# load a presentation
pptx_fname = "old_yg_pptx_template.pptx"

prs = Presentation(pptx_fname)

total_info = xmltodict.parse(prs._element.xml)
try:
    default_clr = total_info['p:presentation']['p:extLst']['p:ext'][0]['p15:sldGuideLst']['p15:guide'][0]['p15:clr']['a:srgbClr']['@val']
except:
    default_clr = total_info['p:presentation']['p:extLst']['p:ext']['p15:sldGuideLst']['p15:guide'][0]['p15:clr']['a:srgbClr']['@val']

title_slide_layout = prs.slide_layouts[3]

theme_info = get_theme_info(title_slide_layout)
theme_colors = theme_info['a:theme']['a:themeElements']['a:clrScheme']
theme_fonts = theme_info['a:theme']['a:themeElements']['a:fontScheme']

for shape in title_slide_layout.shapes:
    print(f"shape_name:\t\t\t{shape.name}")

    # if shape.name == 'Picture Placeholder 6': # Insert Client logo here
    if shape.name == 'Title 14': # Further information
        print('')
    try:
        text_frame = shape.text_frame
        paragraph = text_frame.paragraphs[0]

        text, font_bold, font_italic, font_name, font_size, font_type, font_color = get_info_from_shape(shape)
        print(f"\ttext:\t\t\t{paragraph.text}")
        print(f"\tfont_bold:\t\t{font_bold}")
        print(f"\tfont_italic:\t{font_italic}")
        print(f"\tfont_name:\t\t{font_name}")
        print(f"\tfont_size:\t\t{font_size}")
        print(f"\tfont_type:\t\t{font_type}")
        print(f"\tfont_color:\t\t{font_color}")
        print(f"\tfont_bold:\t\t{font_bold}")
        print('\n-------------------------------------')
    except:
        print('\n-------------------------------------')
        continue

prs.save('old_yg_pptx_template_re_re.pptx')